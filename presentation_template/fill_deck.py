"""Fill the Redrob idea-submission template with our project content.
Reads the blank template, writes a filled copy alongside it. Body text matches the
template style (Manrope SemiBold, ~12pt, #202729). Team name/leader left for the human.
"""
import copy
import os

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

HERE = r"D:/AI_RESUME_BUILDER/presentation_template"
SRC = os.path.join(HERE, "Idea Submission Template _ Redrob.pptx")
OUT = os.path.join(HERE, "Idea Submission - Recruiter-Brain Ranker.pptx")

INK = RGBColor(0x20, 0x27, 0x29)
FONT = "Manrope SemiBold"

# slide index (0-based) -> list of (text, is_header) lines for the body box
BODY = {
    1: [  # Solution Overview
        ("One frozen, human-reviewed JD rubric (rubric.yaml) is the single source of truth — it drives ranking, honeypot defense, AND the reasoning, all from the same understanding of the role.", False),
        ("What differentiates us: we encode the JD's NEGATIVES (disqualifiers / anti-patterns), not just positive keywords. Each maps to a concrete detectable pattern, not a vibe.", False),
        ("Measured on the real pool: 8,545 candidates carry 3+ relevant skills, but only 755 have a corroborating title — a ~10:1 keyword-stuffer ratio. \"Has the skills\" is not discriminative; career coherence is.", False),
    ],
    2: [  # JD Understanding & Candidate Evaluation
        ("Key JD requirements extracted: 5–9 yrs; production embeddings/retrieval; vector-DB & hybrid search; ranking-evaluation literacy (NDCG/MRR); strong Python; shipped at scale; product (not services) company.", False),
        ("Most important signal: domain COHERENCE — title ↔ career ↔ skills must corroborate each other. Weighted far above raw skill tags (title-match 3.0 vs skill-tag 0.4); it's the highest lever on NDCG@10.", False),
        ("Beyond keyword matching: TF-IDF semantic similarity to the JD, behavioral availability (responsive, active, open-to-work), and explicit disqualifier detection for stuffers / consulting-only / honeypots.", False),
    ],
    3: [  # Ranking Methodology
        ("Retrieve/score/rank: offline we extract 42 raw signals + a TF-IDF semantic score into features.parquet; at rank time we apply rubric weights, sort, and take the top 100. No ML model is loaded at rank time.", False),
        ("Scoring formula: final = base_fit (domain coherence + must-haves + semantic) × behavioral_availability − disqualifier_penalties; honeypot_hard → hard floor.", False),
        ("Signals are combined via human-set rubric weights (not a black-box model), so every score is explainable. Tuning re-runs only the ~1 s scorer, never the 90 s feature extraction.", False),
    ],
    4: [  # Explainability & Data Validation
        ("Decisions explained: reasoning is deterministic at rank time (no LLM, no network). Every clause traces to a profile column; skills are named only from the candidate's own list; gaps are stated honestly. 100/100 reasons unique.", False),
        ("Anti-hallucination: a self-check assertion in reasoning.py rejects any clause without a backing data column — unsupported justifications cannot be emitted.", False),
        ("Suspicious profiles: stuffers & consulting-only are soft-penalized (pulled out of the top 10); honeypots (role-months > total career, or \"expert\" with ~0 months used) are hard-floored. Top-100 honeypot rate = 0% (gate exits non-zero if it ever exceeds 10%).", False),
    ],
    5: [  # End-to-End Workflow
        ("OFFLINE (network OK): JD --LLM--> rubric.yaml (human-reviewed, frozen).", False),
        ("BUILD: candidates.jsonl --> features.parquet  (42 raw signals + TF-IDF semantic_sim).", False),
        ("RANK TIME (≤5 min, CPU, no net): load rubric → weighted score → sort → top 100 → generate fact-grounded reason → submission.csv.", False),
        ("One command reproduces everything: python rank.py --candidates ./candidates.jsonl --out ./submission.csv", False),
    ],
    7: [  # Results & Performance
        ("Ranking quality: subset NDCG@10 = 0.9968 (NDCG@50 0.9968, MAP 0.9667) on our hand-labeled eval set — and FLAT across every weight knob, i.e. saturated and robust, not overfit.", False),
        ("Top-100 health: 0% honeypots (pool contains 40), 100/100 domain-corroborated, YoE median 6.6 (inside the 5–9 band).", False),
        ("Runtime/compute: full from-scratch rebuild 109 s (<< 300 s budget), CPU-only, no network; cached path ~1 s. Output is deterministic — reruns are byte-identical.", False),
    ],
    8: [  # Technologies Used
        ("Python, pandas, numpy, scikit-learn (TF-IDF), pyarrow, PyYAML — lightweight, CPU-only, fully reproducible. No heavy model at rank time, which protects the 5-minute budget.", False),
        ("TF-IDF semantic layer chosen over a neural embedding (BGE) for zero heavy deps and full determinism; the architecture keeps a drop-in upgrade path to BGE/E5 if needed.", False),
        ("Streamlit for the interactive sandbox demo. An LLM is used OFFLINE only — to draft the rubric from the JD — then every line is human-reviewed and frozen into version control.", False),
    ],
    9: [  # Submission Assets
        ("GitHub (public, genuine phase-by-phase commit history): github.com/Yash1309-coder/redrobs-hackathon-project-", False),
        ("Live sandbox: Streamlit Cloud demo (URL in submission_metadata.yaml) — interactive top-100 with reasoning + defense funnel.", False),
        ("submission.csv (top 100, candidate_id/rank/score/reasoning), deck.pdf, README with the single reproduce command.", False),
        ("Demo video: [add link]", False),
    ],
}

ARCH = [  # slide 6 (index 5) System Architecture — no body box exists, we add one
    ("Intelligence lives offline; rank time is fast and model-free.", False),
    ("OFFLINE (network OK, no clock):  JD --LLM--> rubric.yaml (frozen)  ·  candidates.jsonl --> features.parquet  [42 raw signals + TF-IDF semantic_sim]", False),
    ("RANK TIME (≤5 min, CPU, no net):  load rubric.yaml  →  score = weighted sum over features  →  sort, take top 100  →  reason from facts  →  submission.csv", False),
    ("Key split: features.parquet holds RAW signals; the scorer applies the rubric WEIGHTS. Re-tuning touches only the ~1 s scorer, never the 90 s extraction — and no model is ever loaded at rank time, so the 5-min budget is never at risk.", False),
]


def fill_box(tf, lines, size=11):
    """Replace a text frame's content with our lines, matching template style."""
    tf.word_wrap = True
    tf.clear()
    for i, (text, _hdr) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = INK


def body_box(slide):
    """The answer box: tallest text box. Titles are ~0.5in tall, bodies ~3.5in."""
    cands = [sh for sh in slide.shapes if sh.has_text_frame and sh.height is not None]
    return max(cands, key=lambda s: s.height) if cands else None


prs = Presentation(SRC)

# Slide 1 — fill Problem Statement only; leave Team Name / Leader for the human.
for sh in prs.slides[0].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith("Problem Statement"):
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = ("  Rank the top 100 of 100,000 candidates against one JD, best-fit first, "
                  "each with an honest fact-grounded reason — on CPU, in under 5 minutes, no network.")
        r.font.name = FONT
        r.font.size = Pt(12)
        r.font.color.rgb = INK

# Content slides with an existing body box.
for idx, lines in BODY.items():
    box = body_box(prs.slides[idx])
    if box is None:
        raise SystemExit(f"no body box on slide index {idx}")
    size = 11 if len(lines) <= 3 else 10
    fill_box(box.text_frame, lines, size=size)

# Slide 7 (index 6) System Architecture — clone a content body box for geometry, add text.
arch_slide = prs.slides[6]
template_box = body_box(prs.slides[4])
new_sp = copy.deepcopy(template_box._element)
arch_slide.shapes._spTree.append(new_sp)
new_box = arch_slide.shapes[-1]
fill_box(new_box.text_frame, ARCH, size=10)

prs.save(OUT)
print("wrote", OUT)
print("slides:", len(prs.slides))
